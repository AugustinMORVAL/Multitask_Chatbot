from abc import ABC, abstractmethod
import io
import fitz
from pptx import Presentation
from PIL import Image
import pytesseract
from openpyxl import load_workbook

class DocumentManager(ABC):
    def convert_to_pdf(self, file):
        pass

class ImageManager(DocumentManager):
    def convert_to_pdf(self, file):
        image = Image.open(io.BytesIO(file.read()))
        text = pytesseract.image_to_string(image)
        return self._text_to_pdf(text)

class TextManager(DocumentManager):
    def convert_to_pdf(self, file):
        text = file.read().decode('utf-8', errors='ignore')
        return self._text_to_pdf(text)

class PowerPointManager(DocumentManager):
    def convert_to_pdf(self, file):
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
        return self._text_to_pdf(text)

class PDFManager(DocumentManager):
    def convert_to_pdf(self, file):
        return io.BytesIO(file.read())

    def _text_to_pdf(self, text):
        pdf = fitz.open()
        page = pdf.new_page()
        page.insert_text((50, 50), text)
        pdf_bytes = pdf.write()
        return io.BytesIO(pdf_bytes)

class ExcelManager(DocumentManager):
    def convert_to_pdf(self, file):
        wb = load_workbook(filename=io.BytesIO(file.read()))
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text += f"Sheet: {sheet}\n"
            for row in ws.iter_rows(values_only=True):
                text += "\t".join(str(cell) for cell in row) + "\n"
            text += "\n"
        return self._text_to_pdf(text)

class DocumentConverterFactory:
    @staticmethod
    def get_converter(file_extension):
        if file_extension in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
            return ImageManager()
        elif file_extension in ['doc', 'docx', 'txt', 'rtf']:
            return TextManager()
        elif file_extension in ['ppt', 'pptx']:
            return PowerPointManager()
        elif file_extension == 'pdf':
            return PDFManager()
        elif file_extension in ['xls', 'xlsx']:
            return ExcelManager()
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
